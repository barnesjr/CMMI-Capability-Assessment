import json
from pathlib import Path
from datetime import datetime
from typing import Optional

try:
    from .models import AssessmentData
except ImportError:
    from models import AssessmentData


MATURITY_LABELS = {
    1: "Initial",
    2: "Managed",
    3: "Defined",
    4: "Quantitatively Managed",
    5: "Optimizing",
}

SCORE_COLORS = {
    1: "ef4444",
    2: "f97316",
    3: "eab308",
    4: "84cc16",
    5: "22c55e",
}

MATURITY_BANDS = [
    {"min": 1.0, "max": 1.5, "label": "Initial", "color": "#ef4444"},
    {"min": 1.5, "max": 2.0, "label": "Managed", "color": "#f97316"},
    {"min": 2.0, "max": 2.5, "label": "Emerging Defined", "color": "#eab308"},
    {"min": 2.5, "max": 3.0, "label": "Defined", "color": "#84cc16"},
    {"min": 3.0, "max": 3.5, "label": "Quantitatively Managed", "color": "#22c55e"},
    {"min": 3.5, "max": 4.0, "label": "Quantitatively Managed+", "color": "#16a34a"},
    {"min": 4.0, "max": 4.5, "label": "Optimizing", "color": "#15803d"},
    {"min": 4.5, "max": 5.1, "label": "Optimizing+", "color": "#166534"},
]


def _timestamp() -> str:
    return datetime.now().strftime("%Y-%m-%d_%H%M%S")


def _get_maturity_band(score: float) -> dict:
    for band in MATURITY_BANDS:
        if band["min"] <= score < band["max"]:
            return band
    if score >= 5.0:
        return MATURITY_BANDS[-1]
    return MATURITY_BANDS[0]


class ExportEngine:
    def __init__(self, base_dir: str, resource_dir: str | None = None):
        self.base_dir = Path(base_dir)
        self.resource_dir = Path(resource_dir) if resource_dir else self.base_dir
        self.exports_dir = self.base_dir / "exports"
        self.templates_dir = self.resource_dir / "templates"

    def _ensure_exports_dir(self):
        self.exports_dir.mkdir(exist_ok=True)

    def _data_dict(self, data: AssessmentData) -> dict:
        return json.loads(data.model_dump_json())

    def _score_label(self, score) -> str:
        if score is None:
            return "Not Scored"
        return MATURITY_LABELS.get(round(score), "N/A")

    def _score_color(self, score) -> str:
        if score is None:
            return "6B6B6B"
        return SCORE_COLORS.get(round(score), "6B6B6B")

    def _avg_score(self, items: list) -> Optional[float]:
        scored = [i for i in items if i.get("score") is not None and not i.get("na", False)]
        if not scored:
            return None
        return sum(i["score"] for i in scored) / len(scored)

    def _pa_score(self, pa: dict) -> Optional[float]:
        items = []
        for ca in pa.get("capability_areas", []):
            items.extend(ca.get("items", []))
        return self._avg_score(items)

    def _composite_score(self, d: dict) -> Optional[float]:
        weights = d.get("scoring_config", {}).get("practice_area_weights", {})
        total_weight = 0.0
        weighted_sum = 0.0
        for group in d.get("category_groups", []):
            for pa in group.get("practice_areas", []):
                score = self._pa_score(pa)
                weight = weights.get(pa["id"], pa.get("weight", 0.05))
                if score is not None:
                    weighted_sum += score * weight
                    total_weight += weight
        if total_weight == 0:
            return None
        return weighted_sum / total_weight

    def _generate_radar_chart(self, data: AssessmentData) -> str:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np

        d = self._data_dict(data)
        labels = []
        scores = []
        targets = []
        target_scores = d.get("target_scores", {})

        for group in d.get("category_groups", []):
            for pa in group.get("practice_areas", []):
                labels.append(pa["name"])
                s = self._pa_score(pa)
                scores.append(s if s is not None else 0)
                targets.append(target_scores.get(pa["id"], 3.0))

        n = len(labels)
        if n == 0:
            fig, ax = plt.subplots(figsize=(6, 6))
            ax.text(0.5, 0.5, "No data", ha="center", va="center")
            self._ensure_exports_dir()
            path = str(self.exports_dir / "radar_chart.png")
            fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
            plt.close(fig)
            return path

        angles = np.linspace(0, 2 * np.pi, n, endpoint=False).tolist()
        scores_plot = scores + [scores[0]]
        targets_plot = targets + [targets[0]]
        angles_plot = angles + [angles[0]]

        fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(polar=True))
        ax.fill(angles_plot, scores_plot, alpha=0.2, color="#3b82f6")
        ax.plot(angles_plot, scores_plot, color="#3b82f6", linewidth=2, label="Current")
        ax.plot(angles_plot, targets_plot, color="#8A8A8E", linewidth=1.5, linestyle="--", label="Target")
        ax.set_xticks(angles)
        ax.set_xticklabels(labels, size=7, wrap=True)
        ax.set_ylim(0, 5)
        ax.set_yticks([1, 2, 3, 4, 5])
        ax.set_yticklabels(["1", "2", "3", "4", "5"], size=7)
        ax.set_title("CMMI Capability Maturity Profile", size=12, pad=20)
        ax.legend(loc="upper right", bbox_to_anchor=(1.3, 1.1), fontsize=8)

        self._ensure_exports_dir()
        path = str(self.exports_dir / "radar_chart.png")
        fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
        plt.close(fig)
        return path

    # ── D-01: Assessment Findings ──────────────────────────────────────

    def export_findings(self, data: AssessmentData) -> str:
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"D-01_Assessment_Findings_{_timestamp()}.docx"
        output_path = self.exports_dir / filename

        from docx import Document
        doc = Document()
        doc.add_heading("CMMI Capability Assessment Findings", 0)

        info = d.get("client_info", {})
        doc.add_paragraph(f"Client: {info.get('name', '')}")
        doc.add_paragraph(f"Industry: {info.get('industry', '')}")
        doc.add_paragraph(f"Date: {info.get('assessment_date', '')}")
        doc.add_paragraph(f"Assessor: {info.get('assessor', '')}")

        composite = self._composite_score(d)
        if composite:
            band = _get_maturity_band(composite)
            doc.add_heading("Overall Maturity", level=1)
            doc.add_paragraph(f"Composite Score: {composite:.2f} — {band['label']}")

        for group in d.get("category_groups", []):
            doc.add_heading(group["name"], level=1)
            for pa in group.get("practice_areas", []):
                pa_s = self._pa_score(pa)
                doc.add_heading(pa["name"], level=2)
                if pa_s:
                    doc.add_paragraph(f"Score: {pa_s:.2f} — {self._score_label(pa_s)}")
                else:
                    doc.add_paragraph("Score: Not yet scored")

                for ca in pa.get("capability_areas", []):
                    ca_score = self._avg_score(ca.get("items", []))
                    doc.add_heading(ca["name"], level=3)
                    doc.add_paragraph(f"Average Score: {ca_score:.2f}" if ca_score else "Not scored")

                    for item in ca.get("items", []):
                        score_val = item.get("score")
                        na = item.get("na", False)
                        if na:
                            text = f"[N/A] {item['text']}"
                        elif score_val:
                            text = f"[{score_val} - {self._score_label(score_val)}] {item['text']}"
                        else:
                            text = f"[--] {item['text']}"
                        if item.get("notes"):
                            text += f"\n  Notes: {item['notes']}"
                        refs = item.get("evidence_references", [])
                        if refs:
                            evidence_str = "; ".join(
                                f"{r.get('document', '')} {r.get('section', '')}"
                                for r in refs if r.get("document")
                            )
                            if evidence_str:
                                text += f"\n  Evidence: {evidence_str}"
                        doc.add_paragraph(text, style="List Bullet")

        doc.save(str(output_path))
        return filename

    # ── D-02: Executive Summary ────────────────────────────────────────

    def export_executive_summary(self, data: AssessmentData) -> str:
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"D-02_Executive_Summary_{_timestamp()}.docx"
        output_path = self.exports_dir / filename
        chart_path = self._generate_radar_chart(data)

        from docx import Document
        from docx.shared import Inches
        doc = Document()
        doc.add_heading("Executive Summary — CMMI Capability Assessment", 0)

        info = d.get("client_info", {})
        doc.add_paragraph(f"Client: {info.get('name', '')}")
        doc.add_paragraph(f"Industry: {info.get('industry', '')}")
        doc.add_paragraph(f"Date: {info.get('assessment_date', '')}")

        composite = self._composite_score(d)
        if composite:
            band = _get_maturity_band(composite)
            doc.add_heading("Overall Maturity", level=1)
            doc.add_paragraph(f"Score: {composite:.2f} — {band['label']}")

        doc.add_heading("Maturity Profile", level=1)
        doc.add_picture(chart_path, width=Inches(5))

        doc.add_heading("Practice Area Scores", level=1)
        for group in d.get("category_groups", []):
            for pa in group.get("practice_areas", []):
                score = self._pa_score(pa)
                if score:
                    doc.add_paragraph(f"{pa['name']}: {score:.2f} — {self._score_label(score)}")
                else:
                    doc.add_paragraph(f"{pa['name']}: N/A")

        # Top gaps
        target_scores = d.get("target_scores", {})
        gaps = []
        for group in d.get("category_groups", []):
            for pa in group.get("practice_areas", []):
                current = self._pa_score(pa)
                target = target_scores.get(pa["id"], 3.0)
                if current is not None:
                    gap = target - current
                    if gap > 0:
                        gaps.append({"name": pa["name"], "current": current, "target": target, "gap": gap})
        gaps.sort(key=lambda g: g["gap"], reverse=True)
        if gaps:
            doc.add_heading("Top Priority Gaps", level=1)
            for g in gaps[:5]:
                doc.add_paragraph(
                    f"{g['name']}: Current {g['current']:.2f} -> Target {g['target']:.1f} (Gap: {g['gap']:.2f})"
                )

        doc.save(str(output_path))
        return filename

    # ── D-03: Gap Analysis & Roadmap ───────────────────────────────────

    def export_gap_analysis(self, data: AssessmentData) -> str:
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"D-03_Gap_Analysis_{_timestamp()}.docx"
        output_path = self.exports_dir / filename
        target_scores = d.get("target_scores", {})

        from docx import Document
        doc = Document()
        doc.add_heading("Gap Analysis & Roadmap", 0)

        info = d.get("client_info", {})
        doc.add_paragraph(f"Client: {info.get('name', '')}")
        doc.add_paragraph(f"Date: {info.get('assessment_date', '')}")

        # Gap matrix table
        doc.add_heading("Gap Matrix", level=1)
        table = doc.add_table(rows=1, cols=5)
        table.style = "Table Grid"
        headers = ["Practice Area", "Current", "Target", "Gap", "Severity"]
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h

        for group in d.get("category_groups", []):
            for pa in group.get("practice_areas", []):
                current = self._pa_score(pa)
                target = target_scores.get(pa["id"], 3.0)
                gap = (target - current) if current else None
                severity = "High" if gap and gap > 1.5 else "Medium" if gap and gap > 0.5 else "Low"
                row = table.add_row().cells
                row[0].text = pa["name"]
                row[1].text = f"{current:.2f}" if current else "N/A"
                row[2].text = f"{target:.1f}"
                row[3].text = f"{gap:.2f}" if gap is not None else "N/A"
                row[4].text = severity

        # Roadmap
        doc.add_heading("Remediation Roadmap", level=1)
        doc.add_heading("30-Day Quick Wins", level=2)
        doc.add_paragraph(
            "Focus on practice areas with scores below 2.0 (Initial/Managed) for immediate improvement opportunities."
        )
        doc.add_heading("60-Day Improvements", level=2)
        doc.add_paragraph(
            "Establish formal, defined processes for practice areas scoring 2.0-3.0."
        )
        doc.add_heading("90-Day Milestones", level=2)
        doc.add_paragraph(
            "Target 'Defined' maturity (3.0+) for critical practice areas with quantitative measurement."
        )
        doc.add_heading("6-12 Month Goals", level=2)
        doc.add_paragraph(
            "Achieve 'Quantitatively Managed' or 'Optimizing' maturity across all practice areas "
            "with continuous improvement cycles and measurable outcomes."
        )

        doc.save(str(output_path))
        return filename

    # ── D-04: Scored Assessment Workbook ───────────────────────────────

    def export_workbook(self, data: AssessmentData) -> str:
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"D-04_Scored_Workbook_{_timestamp()}.xlsx"
        output_path = self.exports_dir / filename

        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment

        wb = Workbook()

        # Dashboard sheet
        ws = wb.active
        ws.title = "Dashboard"
        ws["A1"] = "CMMI Capability Assessment Dashboard"
        ws["A1"].font = Font(bold=True, size=14)
        ws["A3"] = "Client:"
        ws["B3"] = d.get("client_info", {}).get("name", "")
        ws["A4"] = "Industry:"
        ws["B4"] = d.get("client_info", {}).get("industry", "")
        ws["A5"] = "Assessment Date:"
        ws["B5"] = d.get("client_info", {}).get("assessment_date", "")
        ws["A6"] = "Assessor:"
        ws["B6"] = d.get("client_info", {}).get("assessor", "")

        composite = self._composite_score(d)
        ws["A8"] = "Overall Maturity Score:"
        ws["B8"] = round(composite, 2) if composite else "N/A"
        if composite:
            ws["C8"] = _get_maturity_band(composite)["label"]

        row = 10
        header_fill = PatternFill(start_color="3b82f6", end_color="3b82f6", fill_type="solid")
        header_font = Font(bold=True, color="FFFFFF")
        for col, h in enumerate(["Category Group", "Practice Area", "Weight", "Score", "Level"], 1):
            cell = ws.cell(row=row, column=col, value=h)
            cell.font = header_font
            cell.fill = header_fill

        for group in d.get("category_groups", []):
            for pa in group.get("practice_areas", []):
                row += 1
                score = self._pa_score(pa)
                ws.cell(row=row, column=1, value=group["name"])
                ws.cell(row=row, column=2, value=pa["name"])
                ws.cell(row=row, column=3, value=f"{pa.get('weight', 0.05) * 100:.0f}%")
                ws.cell(row=row, column=4, value=round(score, 2) if score else "N/A")
                if score:
                    ws.cell(row=row, column=5, value=self._score_label(score))

        # Per-practice-area sheets
        for group in d.get("category_groups", []):
            for pa in group.get("practice_areas", []):
                sheet_name = pa["name"][:31]
                ws = wb.create_sheet(title=sheet_name)
                ws["A1"] = pa["name"]
                ws["A1"].font = Font(bold=True, size=12)
                ws["A2"] = f"Group: {group['name']}"
                ws["A2"].font = Font(italic=True, color="666666")

                row = 4
                headers = ["Capability Area", "Item #", "Assessment Item", "Score", "Level", "Confidence", "Evidence", "Notes"]
                for col, h in enumerate(headers, 1):
                    cell = ws.cell(row=row, column=col, value=h)
                    cell.font = Font(bold=True)
                    cell.fill = PatternFill(start_color="E2E8F0", end_color="E2E8F0", fill_type="solid")

                for ca in pa.get("capability_areas", []):
                    for item in ca.get("items", []):
                        row += 1
                        score_val = item.get("score")
                        ws.cell(row=row, column=1, value=ca["name"])
                        ws.cell(row=row, column=2, value=item["id"])
                        ws.cell(row=row, column=3, value=item["text"])
                        if item.get("na", False):
                            ws.cell(row=row, column=4, value="N/A")
                            ws.cell(row=row, column=5, value="N/A")
                        else:
                            ws.cell(row=row, column=4, value=score_val)
                            ws.cell(row=row, column=5, value=self._score_label(score_val))
                        ws.cell(row=row, column=6, value=item.get("confidence", ""))
                        refs = item.get("evidence_references", [])
                        evidence_str = "; ".join(
                            f"{r.get('document', '')} {r.get('section', '')}"
                            for r in refs if r.get("document")
                        ) if refs else ""
                        ws.cell(row=row, column=7, value=evidence_str)
                        ws.cell(row=row, column=8, value=item.get("notes", ""))

                        # Color-code score cell
                        if score_val and not item.get("na", False):
                            color = self._score_color(score_val)
                            ws.cell(row=row, column=4).fill = PatternFill(
                                start_color=color, end_color=color, fill_type="solid"
                            )
                            if score_val <= 2:
                                ws.cell(row=row, column=4).font = Font(color="FFFFFF", bold=True)
                            else:
                                ws.cell(row=row, column=4).font = Font(bold=True)

                # Auto-width columns
                for col in ws.columns:
                    max_len = 0
                    col_letter = col[0].column_letter
                    for cell in col:
                        if cell.value:
                            max_len = max(max_len, len(str(cell.value)))
                    ws.column_dimensions[col_letter].width = min(max_len + 2, 60)

        # SVC extension sheet (if enabled)
        if d.get("svc_enabled") and d.get("svc_extension"):
            ws = wb.create_sheet(title="CMMI-SVC")
            ws["A1"] = "CMMI-SVC Extension"
            ws["A1"].font = Font(bold=True, size=12)
            row = 3
            headers = ["Section", "Capability Area", "Item #", "Assessment Item", "Score", "Level", "Confidence", "Notes"]
            for col, h in enumerate(headers, 1):
                cell = ws.cell(row=row, column=col, value=h)
                cell.font = Font(bold=True)
                cell.fill = PatternFill(start_color="E2E8F0", end_color="E2E8F0", fill_type="solid")

            for section in d["svc_extension"].get("sections", []):
                for ca in section.get("capability_areas", []):
                    for item in ca.get("items", []):
                        row += 1
                        score_val = item.get("score")
                        ws.cell(row=row, column=1, value=section["name"])
                        ws.cell(row=row, column=2, value=ca["name"])
                        ws.cell(row=row, column=3, value=item["id"])
                        ws.cell(row=row, column=4, value=item["text"])
                        if item.get("na", False):
                            ws.cell(row=row, column=5, value="N/A")
                            ws.cell(row=row, column=6, value="N/A")
                        else:
                            ws.cell(row=row, column=5, value=score_val)
                            ws.cell(row=row, column=6, value=self._score_label(score_val))
                        ws.cell(row=row, column=7, value=item.get("confidence", ""))
                        ws.cell(row=row, column=8, value=item.get("notes", ""))

        wb.save(str(output_path))
        return filename

    # ── D-05: Out-Brief Presentation ───────────────────────────────────

    def export_outbrief(self, data: AssessmentData) -> str:
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"D-05_OutBrief_{_timestamp()}.pptx"
        output_path = self.exports_dir / filename
        chart_path = self._generate_radar_chart(data)

        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "CMMI Capability Assessment Out-Brief"
        subtitle = next((p for p in slide.placeholders if p.placeholder_format.idx == 1), None)
        if subtitle:
            subtitle.text = d.get("client_info", {}).get("name", "")

        # Overview slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Assessment Overview"
        composite = self._composite_score(d)
        body = next((p for p in slide.placeholders if p.placeholder_format.idx == 1), None)
        if body:
            tf = body.text_frame
            tf.text = f"Client: {d.get('client_info', {}).get('name', '')}"
            tf.add_paragraph().text = f"Industry: {d.get('client_info', {}).get('industry', '')}"
            tf.add_paragraph().text = f"Date: {d.get('client_info', {}).get('assessment_date', '')}"
            tf.add_paragraph().text = f"Assessor: {d.get('client_info', {}).get('assessor', '')}"
            if composite:
                band = _get_maturity_band(composite)
                tf.add_paragraph().text = f"Overall Score: {composite:.2f} — {band['label']}"

        # Radar chart slide
        blank_idx = min(5, len(prs.slide_layouts) - 1)
        slide = prs.slides.add_slide(prs.slide_layouts[blank_idx])
        slide.shapes.add_picture(chart_path, Inches(3), Inches(0.5), Inches(7), Inches(6.5))

        # Per-category-group slides
        for group in d.get("category_groups", []):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = group["name"]
            body = next((p for p in slide.placeholders if p.placeholder_format.idx == 1), None)
            if body:
                tf = body.text_frame
                tf.text = ""
                for pa in group.get("practice_areas", []):
                    score = self._pa_score(pa)
                    p = tf.add_paragraph()
                    if score:
                        p.text = f"{pa['name']}: {score:.2f} — {self._score_label(score)}"
                    else:
                        p.text = f"{pa['name']}: Not scored"

        prs.save(str(output_path))
        return filename

    # ── D-06: Maturity Heatmap ─────────────────────────────────────────

    def export_heatmap(self, data: AssessmentData) -> str:
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"D-06_Maturity_Heatmap_{_timestamp()}.xlsx"
        output_path = self.exports_dir / filename

        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment

        wb = Workbook()
        ws = wb.active
        ws.title = "Heatmap"

        ws["A1"] = "CMMI Capability Maturity Heatmap"
        ws["A1"].font = Font(bold=True, size=14)

        score_colors = {
            1: "EF4444", 2: "F97316", 3: "EAB308", 4: "84CC16", 5: "22C55E",
        }

        row = 3
        ws.cell(row=row, column=1, value="Practice Area").font = Font(bold=True)

        # Find max CA count across all practice areas
        max_cas = 0
        for group in d.get("category_groups", []):
            for pa in group.get("practice_areas", []):
                max_cas = max(max_cas, len(pa.get("capability_areas", [])))

        for i in range(max_cas):
            ws.cell(row=row, column=i + 2, value=f"CA {i + 1}").font = Font(bold=True)
        ws.cell(row=row, column=max_cas + 2, value="PA Avg").font = Font(bold=True)

        for group in d.get("category_groups", []):
            for pa in group.get("practice_areas", []):
                row += 1
                ws.cell(row=row, column=1, value=pa["name"])
                all_items = []
                for ca_idx, ca in enumerate(pa.get("capability_areas", [])):
                    all_items.extend(ca.get("items", []))
                    ca_score = self._avg_score(ca.get("items", []))
                    cell = ws.cell(row=row, column=ca_idx + 2)
                    if ca_score is not None:
                        cell.value = round(ca_score, 2)
                        rounded = min(5, max(1, round(ca_score)))
                        color = score_colors.get(rounded, "FFFFFF")
                        cell.fill = PatternFill(start_color=color, end_color=color, fill_type="solid")
                        if rounded <= 2:
                            cell.font = Font(color="FFFFFF", bold=True)
                        else:
                            cell.font = Font(bold=True)
                    else:
                        cell.value = "--"
                pa_avg = self._avg_score(all_items)
                avg_cell = ws.cell(row=row, column=max_cas + 2)
                if pa_avg is not None:
                    avg_cell.value = round(pa_avg, 2)
                    avg_cell.font = Font(bold=True)

        # Legend
        row += 2
        ws.cell(row=row, column=1, value="Legend:").font = Font(bold=True)
        for score_val, label in MATURITY_LABELS.items():
            row += 1
            cell = ws.cell(row=row, column=1, value=f"{score_val} — {label}")
            color = score_colors[score_val]
            cell.fill = PatternFill(start_color=color, end_color=color, fill_type="solid")
            if score_val <= 2:
                cell.font = Font(color="FFFFFF")

        # Auto-width
        ws.column_dimensions["A"].width = 35
        for i in range(max_cas + 2):
            col_letter = chr(66 + i) if (66 + i) <= 90 else None
            if col_letter:
                ws.column_dimensions[col_letter].width = 12

        wb.save(str(output_path))
        return filename

    # ── D-07: Quick Wins Report ────────────────────────────────────────

    def export_quick_wins(self, data: AssessmentData) -> str:
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"D-07_Quick_Wins_{_timestamp()}.docx"
        output_path = self.exports_dir / filename

        from docx import Document
        doc = Document()
        doc.add_heading("CMMI Quick Wins Report", 0)

        info = d.get("client_info", {})
        doc.add_paragraph(f"Client: {info.get('name', '')}")
        doc.add_paragraph(f"Date: {info.get('assessment_date', '')}")

        doc.add_heading("Methodology", level=1)
        doc.add_paragraph(
            "Quick wins are identified as assessment items scored 1-2 (Initial or Managed) "
            "that have the highest potential for improvement. Items are prioritized by practice area weight "
            "multiplied by the score gap to target."
        )

        # Collect quick win candidates
        target_scores = d.get("target_scores", {})
        weights = d.get("scoring_config", {}).get("practice_area_weights", {})
        candidates = []
        for group in d.get("category_groups", []):
            for pa in group.get("practice_areas", []):
                weight = weights.get(pa["id"], pa.get("weight", 0.05))
                target = target_scores.get(pa["id"], 3.0)
                for ca in pa.get("capability_areas", []):
                    for item in ca.get("items", []):
                        score = item.get("score")
                        if score is not None and score <= 2 and not item.get("na", False):
                            gap = target - score
                            priority = weight * gap
                            candidates.append({
                                "practice_area": pa["name"],
                                "ca": ca["name"],
                                "item_id": item["id"],
                                "text": item["text"],
                                "score": score,
                                "target": target,
                                "gap": gap,
                                "priority": priority,
                                "notes": item.get("notes", ""),
                            })

        candidates.sort(key=lambda c: c["priority"], reverse=True)

        if not candidates:
            doc.add_heading("Results", level=1)
            doc.add_paragraph("No quick win candidates found. All items are scored 3 or above, or are unscored.")
        else:
            from collections import defaultdict
            by_pa: dict[str, list] = defaultdict(list)
            for c in candidates[:30]:  # Top 30
                by_pa[c["practice_area"]].append(c)

            for pa_name, items in by_pa.items():
                doc.add_heading(pa_name, level=1)
                table = doc.add_table(rows=1, cols=4)
                table.style = "Table Grid"
                for i, h in enumerate(["Item", "Current", "Target", "Gap"]):
                    table.rows[0].cells[i].text = h
                for item in items:
                    row = table.add_row().cells
                    row[0].text = f"[{item['item_id']}] {item['text'][:80]}"
                    row[1].text = f"{item['score']} ({self._score_label(item['score'])})"
                    row[2].text = f"{item['target']:.1f}"
                    row[3].text = f"{item['gap']:.1f}"

        doc.save(str(output_path))
        return filename

    # ── D-08: CMMI Roadmap ─────────────────────────────────────────────

    def export_cmmi_roadmap(self, data: AssessmentData) -> str:
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"D-08_CMMI_Roadmap_{_timestamp()}.docx"
        output_path = self.exports_dir / filename

        from docx import Document
        doc = Document()
        doc.add_heading("CMMI Capability Maturity Roadmap", 0)

        info = d.get("client_info", {})
        doc.add_paragraph(f"Client: {info.get('name', '')}")
        doc.add_paragraph(f"Date: {info.get('assessment_date', '')}")

        composite = self._composite_score(d)
        if composite:
            doc.add_heading("Current State", level=1)
            doc.add_paragraph(f"Overall Maturity Score: {composite:.2f} — {self._score_label(composite)}")

        # Collect PA scores for roadmap planning
        pa_scores = []
        for group in d.get("category_groups", []):
            for pa in group.get("practice_areas", []):
                score = self._pa_score(pa)
                pa_scores.append({"name": pa["name"], "group": group["name"], "score": score})

        # Level-by-level progression plan
        for target_level in range(2, 6):
            label = MATURITY_LABELS.get(target_level, "")
            doc.add_heading(f"Level {target_level}: {label}", level=1)

            doc.add_heading("Objective", level=2)
            objectives = {
                2: "Establish basic management practices. Projects are planned, performed, and controlled. "
                   "Existing practices are retained during times of stress.",
                3: "Processes are well characterized, understood, and described in standards, procedures, tools, and methods. "
                   "Organizational standard processes are established and improved over time.",
                4: "Quantitative objectives for quality and process performance are established and used as criteria. "
                   "Quantitative objectives are based on customer, end-user, organization, and process needs.",
                5: "Focus on continually improving process performance through incremental and innovative improvements. "
                   "Quality and process performance objectives are established, continually revised to reflect changing business objectives.",
            }
            doc.add_paragraph(objectives.get(target_level, ""))

            # Practice areas below this level
            below = [p for p in pa_scores if p["score"] is not None and p["score"] < target_level]
            at_or_above = [p for p in pa_scores if p["score"] is not None and p["score"] >= target_level]
            unscored = [p for p in pa_scores if p["score"] is None]

            doc.add_heading("Practice Areas Requiring Improvement", level=2)
            if below:
                table = doc.add_table(rows=1, cols=3)
                table.style = "Table Grid"
                for i, h in enumerate(["Practice Area", "Current Score", "Gap to Level"]):
                    table.rows[0].cells[i].text = h
                for p in sorted(below, key=lambda x: x["score"]):
                    row = table.add_row().cells
                    row[0].text = p["name"]
                    row[1].text = f"{p['score']:.2f}" if p["score"] else "N/A"
                    row[2].text = f"{target_level - p['score']:.2f}" if p["score"] else "N/A"
            else:
                doc.add_paragraph("All scored practice areas meet or exceed this level.")

            doc.add_heading("Practice Areas Meeting This Level", level=2)
            if at_or_above:
                for p in at_or_above:
                    doc.add_paragraph(f"{p['name']}: {p['score']:.2f}", style="List Bullet")
            else:
                doc.add_paragraph("No practice areas currently meet this level.")

        doc.save(str(output_path))
        return filename

    # ── D-09: SVC Alignment Report ─────────────────────────────────────

    def export_svc_alignment(self, data: AssessmentData) -> str:
        self._ensure_exports_dir()
        d = self._data_dict(data)
        filename = f"D-09_SVC_Alignment_{_timestamp()}.docx"
        output_path = self.exports_dir / filename

        from docx import Document
        doc = Document()
        doc.add_heading("CMMI-SVC Alignment Report", 0)

        info = d.get("client_info", {})
        doc.add_paragraph(f"Client: {info.get('name', '')}")
        doc.add_paragraph(f"Date: {info.get('assessment_date', '')}")

        if not d.get("svc_enabled") or not d.get("svc_extension"):
            doc.add_paragraph("CMMI-SVC extension is not enabled for this assessment.")
            doc.save(str(output_path))
            return filename

        extension = d["svc_extension"]
        all_svc_items = []
        for section in extension.get("sections", []):
            for ca in section.get("capability_areas", []):
                all_svc_items.extend(ca.get("items", []))

        overall = self._avg_score(all_svc_items)
        if overall:
            doc.add_heading("Overall SVC Alignment", level=1)
            doc.add_paragraph(f"Score: {overall:.2f} — {self._score_label(overall)}")

        for section in extension.get("sections", []):
            section_items = []
            for ca in section.get("capability_areas", []):
                section_items.extend(ca.get("items", []))
            score = self._avg_score(section_items)
            doc.add_heading(section["name"], level=1)
            if score:
                doc.add_paragraph(f"Score: {score:.2f} — {self._score_label(score)}")
            else:
                doc.add_paragraph("Not yet scored")

            for ca in section.get("capability_areas", []):
                ca_score = self._avg_score(ca.get("items", []))
                doc.add_heading(ca["name"], level=2)
                doc.add_paragraph(f"Average: {ca_score:.2f}" if ca_score else "Not scored")
                for item in ca.get("items", []):
                    s = item.get("score")
                    if item.get("na"):
                        text = f"[N/A] {item['text']}"
                    elif s:
                        text = f"[{s} - {self._score_label(s)}] {item['text']}"
                    else:
                        text = f"[--] {item['text']}"
                    if item.get("notes"):
                        text += f" — {item['notes']}"
                    doc.add_paragraph(text, style="List Bullet")

        doc.save(str(output_path))
        return filename

    # ── Export All ─────────────────────────────────────────────────────

    def export_all(self, data: AssessmentData) -> list[str]:
        filenames = [
            self.export_findings(data),
            self.export_executive_summary(data),
            self.export_gap_analysis(data),
            self.export_workbook(data),
            self.export_outbrief(data),
            self.export_heatmap(data),
            self.export_quick_wins(data),
            self.export_cmmi_roadmap(data),
        ]
        if data.svc_enabled:
            filenames.append(self.export_svc_alignment(data))
        return filenames
